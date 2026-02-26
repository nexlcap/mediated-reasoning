"""Document/PDF loader for grounding context (Item 29).

Supports:
  - PDF (via pypdf, falls back to pdfminer, then pytesseract OCR)
  - Plain text (.txt, .md)
  - Word (.docx, via python-docx)
  - PowerPoint (.pptx, via python-pptx)
  - Excel (.xlsx, via openpyxl)
  - Raises DocumentLoadError with a user-friendly message on failure.
"""
from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Optional

from src.models.schemas import DocumentContext

logger = logging.getLogger(__name__)

# ── constants ─────────────────────────────────────────────────────────────────

MAX_CHARS = 120_000  # ~30 k tokens — keeps context window sane
TRUNCATION_NOTE = "\n\n[Document truncated to fit context window.]"


class DocumentLoadError(Exception):
    """Raised when a document cannot be loaded or parsed."""


# ── helpers ───────────────────────────────────────────────────────────────────


def _truncate(text: str) -> str:
    if len(text) <= MAX_CHARS:
        return text
    return text[:MAX_CHARS] + TRUNCATION_NOTE


def _extract_pdf_pypdf(data: bytes) -> tuple[str, int]:
    """Extract text via pypdf (preferred — pure-Python, no binary deps)."""
    try:
        import pypdf  # type: ignore
    except ImportError:
        raise ImportError("pypdf")

    reader = pypdf.PdfReader(io.BytesIO(data))
    pages = []
    for page in reader.pages:
        text = page.extract_text() or ""
        pages.append(text)
    return "\n\n".join(pages), len(reader.pages)


def _extract_pdf_pdfminer(data: bytes) -> tuple[str, int]:
    """Fallback extraction via pdfminer.six."""
    try:
        from pdfminer.high_level import extract_text as pm_extract  # type: ignore
        from pdfminer.pdfpage import PDFPage  # type: ignore
    except ImportError:
        raise ImportError("pdfminer.six")

    text = pm_extract(io.BytesIO(data))
    page_count = sum(1 for _ in PDFPage.get_pages(io.BytesIO(data)))
    return text or "", page_count


def _extract_pdf_ocr(data: bytes) -> tuple[str, int]:
    """Last-resort OCR via pdf2image + pytesseract (requires system deps)."""
    try:
        import pytesseract  # type: ignore
        from pdf2image import convert_from_bytes  # type: ignore
    except ImportError:
        raise ImportError("pdf2image and pytesseract")

    images = convert_from_bytes(data)
    pages = [pytesseract.image_to_string(img) for img in images]
    return "\n\n".join(pages), len(images)


def _extract_pdf(data: bytes) -> tuple[str, int, str]:
    """Return (text, page_count, method) trying extractors in priority order."""
    for extractor, label in [
        (_extract_pdf_pypdf, "text"),
        (_extract_pdf_pdfminer, "text"),
        (_extract_pdf_ocr, "ocr"),
    ]:
        try:
            text, pages = extractor(data)
            if text.strip():
                return text, pages, label
        except ImportError as exc:
            logger.debug("PDF extractor unavailable (%s): %s", label, exc)
        except Exception as exc:  # noqa: BLE001
            logger.warning("PDF extractor failed (%s): %s", label, exc)

    raise DocumentLoadError(
        "Could not extract text from PDF. "
        "Install pypdf (`pip install pypdf`) or pdfminer.six for text-based PDFs; "
        "install pdf2image + pytesseract for scanned/image PDFs."
    )


# ── Office format extractors ──────────────────────────────────────────────────


def _extract_docx(data: bytes) -> str:
    """Extract text from a Word .docx file via python-docx."""
    try:
        import docx  # type: ignore
    except ImportError:
        raise DocumentLoadError(
            "python-docx is required for Word files: pip install python-docx"
        )
    doc = docx.Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append("\t".join(cells))
    return "\n\n".join(paragraphs)


def _extract_pptx(data: bytes) -> tuple[str, int]:
    """Extract text from a PowerPoint .pptx file via python-pptx."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        raise DocumentLoadError(
            "python-pptx is required for PowerPoint files: pip install python-pptx"
        )
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides), len(prs.slides)


def _extract_xlsx(data: bytes) -> tuple[str, int]:
    """Extract cell data from an Excel .xlsx file via openpyxl."""
    try:
        import openpyxl  # type: ignore
    except ImportError:
        raise DocumentLoadError(
            "openpyxl is required for Excel files: pip install openpyxl"
        )
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append("\t".join(cells))
        if rows:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    return "\n\n".join(sheets), len(wb.sheetnames)


# ── public API ────────────────────────────────────────────────────────────────


def load_document(source: str | Path | bytes, filename: Optional[str] = None) -> DocumentContext:
    """Load a document from a file path or raw bytes.

    Parameters
    ----------
    source:
        A filesystem path (str or Path) OR raw bytes of the file content.
    filename:
        Original filename — used to determine format when *source* is bytes.
        Required when source is bytes.

    Returns
    -------
    DocumentContext with extracted text, page count, and extraction method.
    """
    if isinstance(source, (str, Path)):
        path = Path(source)
        if not path.exists():
            raise DocumentLoadError(f"File not found: {path}")
        if not path.is_file():
            raise DocumentLoadError(f"Not a file: {path}")
        data = path.read_bytes()
        fname = filename or path.name
        suffix = path.suffix.lower()
    else:
        # raw bytes
        if not filename:
            raise DocumentLoadError("filename is required when source is bytes")
        data = source
        fname = filename
        suffix = Path(filename).suffix.lower()

    if suffix == ".pdf":
        text, page_count, method = _extract_pdf(data)
    elif suffix in {".txt", ".md", ".rst", ".text"}:
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            text = data.decode("latin-1")
        page_count = None
        method = "text"
    elif suffix == ".docx":
        text = _extract_docx(data)
        page_count = None
        method = "text"
    elif suffix == ".pptx":
        text, page_count = _extract_pptx(data)
        method = "text"
    elif suffix in {".xlsx", ".xls"}:
        text, page_count = _extract_xlsx(data)
        method = "text"
    else:
        raise DocumentLoadError(
            f"Unsupported file type: '{suffix}'. "
            "Supported formats: .pdf, .txt, .md, .rst, .docx, .pptx, .xlsx"
        )

    text = _truncate(text.strip())

    if not text:
        raise DocumentLoadError(
            f"No text could be extracted from '{fname}'. "
            "The file may be empty or corrupt."
        )

    return DocumentContext(
        filename=fname,
        content=text,
        page_count=page_count,
        extraction_method=method,
    )
